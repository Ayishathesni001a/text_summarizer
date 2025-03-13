from pptx import Presentation
from pptx.util import Inches

class PPTGenerator:
    def __init__(self):
        self.presentation = Presentation()
    
    def create_presentation(self, title, content):
        # Create title slide
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        # Create content slides
        for i, slide_content in enumerate(content, 1):
            slide_layout = self.presentation.slide_layouts[1]
            slide = self.presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = f"Slide {i}"
            content_placeholder.text = slide_content
        
        # Save the presentation
        self.presentation.save("summary_presentation.pptx")
        return "summary_presentation.pptx"
    
    def add_slide(self, title, content):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
    
    def save_presentation(self, file_path):
        self.presentation.save(file_path)
        return file_path