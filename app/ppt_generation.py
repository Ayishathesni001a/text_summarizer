from pptx import Presentation
from pptx.util import Inches
import os
import re

class PPTGenerator:
    def __init__(self):
        self.presentation = Presentation()
        self.output_dir = "output_ppt"
        
        # Create output directory if it doesn't exist
        os.makedirs(self.output_dir, exist_ok=True)

    def create_from_text_file(self, input_file_path):
        """Create presentation from a text file in output_summary folder"""
        try:
            with open(input_file_path, 'r') as file:
                content = file.read()
            
            # Extract title from first line
            title = os.path.basename(input_file_path).replace('.txt', '')
            
            # Split content into slides using a pattern like "Slide 1:"
            slides_content = re.split(r'Slide \d+:', content)[1:]  # Skip first empty element
            
            # Create presentation
            output_path = os.path.join(self.output_dir, f"{title}.pptx")
            self._create_presentation(title, slides_content)
            self.save_presentation(output_path)
            
            return output_path
            
        except Exception as e:
            print(f"Error processing {input_file_path}: {str(e)}")
            return None

    def _create_presentation(self, title, slides_content):
        """Internal method to create presentation structure"""
        # Title slide
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title_slide.shapes.title.text = title

        # Content slides
        for idx, content in enumerate(slides_content, 1):
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = f"{title} - Part {idx}"
            content_shape.text = content.strip()

    def save_presentation(self, file_path):
        """Save presentation to specified path"""
        self.presentation.save(file_path)
        return file_path

def process_summary_folder(input_dir="output_summary"):
    generator = PPTGenerator()
    
    # Process all text files in input directory
    for filename in os.listdir(input_dir):
        if filename.endswith(".txt"):
            input_path = os.path.join(input_dir, filename)
            output_path = generator.create_from_text_file(input_path)
            
            if output_path:
                print(f"Created presentation: {output_path}")
    
    return generator.output_dir

# Example usage
if __name__ == "__main__":
    output_directory = process_summary_folder()
    print(f"All presentations saved to: {output_directory}")